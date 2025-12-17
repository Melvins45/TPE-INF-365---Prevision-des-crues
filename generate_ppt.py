from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Prévision des Crues"
    subtitle.text = "Projet TPE INF 365 - Groupe 24\nUniversité de Douala - Faculté des Sciences\nExaminateur : Dr Justin MOSKOLAI"

    # Slide 1.5: Team Members
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Membres du Groupe 24"
    content = slide.placeholders[1]
    content.text = (
        "• NITOPOP JEATSA GUILLAUME MELVIN (CHEF) - 20S43003\n"
        "• DEMANOU KEMKENG DILAN - 25S03516\n"
        "• NOSSI YIMGO LYNDSEY SULIVANE - 23S87713\n"
        "• TCHIEUTCHOUA FOTEPING ASHLEY MEGANE - 23S87863\n"
        "• WANGA POUYA KAVEN SAMIRA - 23S88070"
    )

    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction & Contexte"
    content = slide.placeholders[1]
    content.text = (
        "Contexte : Prévention des inondations à Douala, Cameroun.\n"
        "Objectif : Prédire la probabilité d'inondation (FloodProbability) en fonction de divers facteurs environnementaux.\n\n"
        "Données : Dataset 'flood.csv' contenant des caractéristiques telles que :\n"
        "- Intensité de la mousson\n"
        "- Topographie et drainage\n"
        "- Déforestation\n"
        "- Urbanisation\n"
        "- Changement climatique"
    )

    # Slide 3: Méthodologie
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Méthodologie"
    content = slide.placeholders[1]
    content.text = (
        "1. Préparation des données :\n"
        "   - Nettoyage : Traitement des valeurs manquantes et doublons.\n"
        "   - Outliers : Suppression ou plafonnement via la méthode IQR (Interquartile Range).\n"
        "   - Normalisation : Utilisation de StandardScaler pour mettre les features à la même échelle.\n\n"
        "2. Modélisation :\n"
        "   - Régression Linéaire (Baseline)\n"
        "   - Forêt Aléatoire (Random Forest)\n"
        "   - Réseau de Neurones (Deep Learning)"
    )

    # Slide 3.5: Analyse Exploratoire (Graphiques)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Analyse Exploratoire"
    content = slide.placeholders[1]
    content.text = (
        "Matrice de Corrélation :\n"
        "Permet de visualiser les relations linéaires entre les variables.\n\n"
        "[ESPACE POUR IMAGE: Matrice de Corrélation (Heatmap)]\n"
        "(Voir notebook: modeleRegressionLineaire.ipynb ou modeleForetAleatoire.ipynb)"
    )

    # Slide 4: Modèles Utilisés
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Détails des Modèles"
    content = slide.placeholders[1]
    content.text = (
        "Régression Linéaire :\n"
        "- Modèle simple pour établir une base de performance.\n\n"
        "Forêt Aléatoire (Random Forest) :\n"
        "- Ensemble de 100 arbres de décision.\n"
        "- Robuste contre le surapprentissage.\n"
        "- Capture les relations non-linéaires.\n\n"
        "Réseau de Neurones Séquentiel :\n"
        "- Architecture : 3 couches denses (64, 32, 1 neurones).\n"
        "- Activation : ReLU pour les couches cachées, Linéaire pour la sortie.\n"
        "- Optimiseur : Adam | Perte : MSE."
    )

    # Slide 4.1: Résultats - Régression Linéaire
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Résultats : Régression Linéaire"
    content = slide.placeholders[1]
    content.text = (
        "Nous commençons par observer les résultats de la Régression Linéaire, notre modèle de base. "
        "Ce graphique met en relation les probabilités d'inondation réelles (axe horizontal) avec celles prédites par le modèle (axe vertical). "
        "La ligne rouge représente la perfection : si tous les points étaient dessus, le modèle aurait 100% de réussite. "
        "Ici, la dispersion des points nous montre que ce modèle simple peine à capturer toute la complexité des facteurs environnementaux, "
        "ce qui entraîne des erreurs de prévision visibles et une fiabilité limitée pour des cas critiques.\n\n"
        "[ESPACE POUR IMAGE: Actual vs Predicted FloodProbability]\n"
        "(Voir notebook: modeleRegressionLineaire.ipynb)"
    )

    # Slide 4.2: Résultats - Forêt Aléatoire
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Résultats : Forêt Aléatoire"
    content = slide.placeholders[1]
    content.text = (
        "Passons ensuite au modèle de Forêt Aléatoire, qui utilise la puissance de 100 arbres de décision travaillant ensemble. "
        "Contrairement au modèle précédent, vous remarquerez que le nuage de points est beaucoup plus resserré autour de la diagonale rouge. "
        "Cela signifie que les prédictions sont nettement plus fiables et précises. Ce modèle parvient mieux à comprendre les interactions "
        "non-linéaires entre les variables (comme l'effet combiné de la pluie et de l'urbanisation), réduisant ainsi considérablement l'incertitude.\n\n"
        "[ESPACE POUR IMAGE: Actual vs Predicted FloodProbability]\n"
        "(Voir notebook: modeleForetAleatoire.ipynb)"
    )

    # Slide 4.3: Résultats - Réseau de Neurones
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Résultats : Réseau de Neurones"
    content = slide.placeholders[1]
    content.text = (
        "Enfin, voici les performances de notre Réseau de Neurones Séquentiel, une approche de Deep Learning avancée. "
        "Ce modèle a appris à pondérer finement chaque caractéristique à travers ses couches de neurones artificiels successives. "
        "Le graphique montre une adéquation quasi-parfaite entre les prédictions et la réalité, avec des points qui épousent très fidèlement la ligne directrice. "
        "C'est le signe que le modèle a réussi à généraliser les règles sous-jacentes des inondations, offrant ainsi la meilleure fiabilité pour notre système d'alerte.\n\n"
        "[ESPACE POUR IMAGE: Actual vs Predicted]\n"
        "(Voir notebook: modeleReseauNeuronesSequential.ipynb)"
    )

    # Slide 5: Validation Croisée (Cross-Validation)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Validation Croisée"
    content = slide.placeholders[1]
    content.text = (
        "Pourquoi ?\n"
        "- Pour évaluer la robustesse des modèles.\n"
        "- Évite que le modèle soit performant uniquement sur une partie spécifique des données.\n\n"
        "Méthode (K-Fold) :\n"
        "- Division des données en 5 parties (folds).\n"
        "- Entraînement sur 4 parties, test sur 1.\n"
        "- Répétition 5 fois et moyenne des scores (MSE, R2)."
    )

    # Slide 6: Résultats (Aperçu)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Résultats Préliminaires"
    content = slide.placeholders[1]
    content.text = (
        "Réseau de Neurones :\n"
        "- MSE (Mean Squared Error) très faible (~1.4e-05) sur les données de validation.\n"
        "- Indique une très bonne capacité prédictive.\n\n"
        "Comparaison :\n"
        "- La validation croisée permettra de confirmer si le Random Forest ou le Réseau de Neurones est le meilleur modèle généralisable.\n"
        "- Le script 'validation_croisee.py' a été créé pour automatiser ce test."
    )

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = (
        "- Le projet a permis d'explorer plusieurs approches de modélisation.\n"
        "- Le prétraitement des données (gestion des outliers, scaling) est crucial.\n"
        "- Les modèles complexes (NN, RF) semblent prometteurs.\n"
        "- Prochaines étapes : Tuning des hyperparamètres et déploiement du meilleur modèle."
    )

    prs.save('presentation_projet.pptx')
    print("Présentation générée : presentation_projet.pptx")

if __name__ == "__main__":
    create_presentation()
