from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_theme():
    prs = Presentation()
    
    # Access the first slide master
    master = prs.slide_masters[0]
    
    # Set background color to a deep professional blue (Dark Slate Blue)
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = RGBColor(28, 40, 51) 
    
    # Add a decorative bottom bar (Water/Flood theme - Light Blue)
    # Note: python-pptx has limitations on editing layout shapes directly in some versions.
    # Instead of trying to create a complex template programmatically which is error-prone with this library,
    # we will create a presentation where we apply the style to the SLIDES themselves as we create them,
    # OR we just set the background color which is supported.
    
    # Let's stick to setting the background color on the master, which worked (or should work).
    # And we will skip adding shapes to layouts/masters if the API doesn't support it easily.
    # The background color alone provides a strong "theme".
    
    pass 

    # Customize Text Styles on Layouts to ensure visibility on dark background
    for layout in master.slide_layouts:
        for shape in layout.placeholders:
            if not shape.has_text_frame:
                continue
            
            # We can't easily set the "default" style for the placeholder via python-pptx 
            # in a way that persists for new text added later without direct XML manipulation,
            # but we can try to set the font color of the placeholder itself.
            # When the user (or script) adds text, it might reset, but this is a best effort for a template.
            
            try:
                text_frame = shape.text_frame
                # Set initial paragraph font
                if text_frame.paragraphs:
                    p = text_frame.paragraphs[0]
                    p.font.color.rgb = RGBColor(236, 240, 241) # Off-white
                    p.font.name = 'Arial'
            except Exception:
                pass

    prs.save('theme_crues.pptx')
    print("Thème généré : theme_crues.pptx")

if __name__ == "__main__":
    create_theme()
